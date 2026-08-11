#!/usr/bin/env python3
"""
Hospital 360 — Markdown-to-PowerPoint Generator
Converts hospital360_demo_deck.md into a styled .pptx file.
"""

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SNOWFLAKE_BLUE = RGBColor(0x29, 0xB5, 0xE8)
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
TABLE_HEADER_BG = RGBColor(0x29, 0xB5, 0xE8)
TABLE_ALT_BG = RGBColor(0xEB, 0xF7, 0xFC)
CODE_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

SRC = Path(__file__).parent / "hospital360_demo_deck.md"
DST = Path(__file__).parent / "hospital360_demo_deck.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_footer(slide, prs):
    """Add a small footer to the bottom of a slide."""
    left = Inches(0.5)
    top = SLIDE_HEIGHT - Inches(0.4)
    width = SLIDE_WIDTH - Inches(1)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Hospital 360  |  Snowflake + Epic Healthcare Analytics"
    p.font.size = Pt(9)
    p.font.color.rgb = MED_GRAY
    p.alignment = PP_ALIGN.CENTER


def set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a simple textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_rich_bullets(slide, left, top, width, height, lines, font_size=14,
                     color=DARK_TEXT):
    """Add bullet points from markdown lines (supports **bold** fragments)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        # Determine indent level
        stripped = line.lstrip()
        indent_level = 0
        if line.startswith("   ") or line.startswith("\t"):
            indent_level = 1

        # Remove leading - or * or numbered prefix
        stripped = re.sub(r"^[\-\*]\s+", "", stripped)
        stripped = re.sub(r"^\d+\.\s+", "", stripped)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = indent_level
        p.space_before = Pt(4)
        p.space_after = Pt(2)

        # Parse **bold** and `code` fragments
        parts = re.split(r"(\*\*.*?\*\*|`[^`]+`)", stripped)
        for part in parts:
            if not part:
                continue
            run = p.add_run()
            if part.startswith("**") and part.endswith("**"):
                run.text = part[2:-2]
                run.font.bold = True
            elif part.startswith("`") and part.endswith("`"):
                run.text = part[1:-1]
                run.font.name = "Consolas"
                run.font.size = Pt(font_size - 1)
            else:
                run.text = part
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = run.font.name or "Calibri"

    return tf


def add_table(slide, left, top, width, rows_data, col_widths=None):
    """Add a styled table from parsed markdown table data."""
    if not rows_data or len(rows_data) < 1:
        return

    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    height = Inches(0.35 * n_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            if i < n_cols:
                table.columns[i].width = w
    else:
        col_w = int(width / n_cols)
        for i in range(n_cols):
            table.columns[i].width = col_w

    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text.strip()

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_TEXT

            # Header row styling
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG

    return table_shape


def add_code_block(slide, left, top, width, code_text, font_size=10):
    """Add a code block with gray background."""
    lines = code_text.strip().split("\n")
    height = Inches(0.22 * len(lines) + 0.3)
    # Background rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()

    # Code text
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                     width - Inches(0.3), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    return height


# ---------------------------------------------------------------------------
# Markdown parsing helpers
# ---------------------------------------------------------------------------
def parse_md_table(lines):
    """Parse markdown table lines into list of row-lists. Skips separator row."""
    rows = []
    for line in lines:
        line = line.strip()
        if not line.startswith("|"):
            continue
        # Skip separator rows like |---|---|
        if re.match(r"^\|[\s\-\|:]+\|$", line):
            continue
        cells = [c.strip() for c in line.split("|")[1:-1]]
        rows.append(cells)
    return rows


def split_slide_content(text):
    """Break a slide's raw markdown into structured elements."""
    elements = []
    lines = text.strip().split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]

        # Code block
        if line.strip().startswith("```"):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing ```
            elements.append(("code", "\n".join(code_lines)))
            continue

        # Table (collect consecutive | lines)
        if line.strip().startswith("|"):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            elements.append(("table", table_lines))
            continue

        # H1
        if line.startswith("# ") and not line.startswith("## "):
            elements.append(("h1", line[2:].strip()))
            i += 1
            continue

        # H2
        if line.startswith("## "):
            elements.append(("h2", line[3:].strip()))
            i += 1
            continue

        # H3
        if line.startswith("### "):
            elements.append(("h3", line[4:].strip()))
            i += 1
            continue

        # Bullet / numbered list
        if re.match(r"^[\-\*]\s+", line.strip()) or re.match(r"^\d+\.\s+", line.strip()):
            bullet_lines = []
            while i < len(lines):
                l = lines[i]
                ls = l.strip()
                if (re.match(r"^[\-\*]\s+", ls) or re.match(r"^\d+\.\s+", ls) or
                        l.startswith("   ") or l.startswith("\t")):
                    bullet_lines.append(l)
                    i += 1
                elif ls == "":
                    i += 1
                    # Check if next line continues the list
                    if i < len(lines) and (re.match(r"^[\-\*]\s+", lines[i].strip()) or
                                           re.match(r"^\d+\.\s+", lines[i].strip())):
                        continue
                    break
                else:
                    break
            elements.append(("bullets", bullet_lines))
            continue

        # Blockquote
        if line.strip().startswith(">"):
            elements.append(("quote", line.strip()[1:].strip()))
            i += 1
            continue

        # Regular text
        text_content = line.strip()
        if text_content and not text_content.startswith("<!--"):
            elements.append(("text", text_content))
        i += 1

    return elements


# ---------------------------------------------------------------------------
# Slide builders for different slide types
# ---------------------------------------------------------------------------
def build_title_slide(prs, elements):
    """Build the title slide (slide 1) with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    h1 = ""
    h2 = ""
    subtitle_text = ""
    tagline = ""
    for etype, content in elements:
        if etype == "h1" and not h1:
            h1 = content
        elif etype == "h2" and not h2:
            h2 = content
        elif etype == "text":
            if "**" in content:
                subtitle_text = content.replace("**", "")
            elif not tagline:
                tagline = content

    # Main title
    add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                h1, font_size=44, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    if h2:
        add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                    h2, font_size=24, color=SNOWFLAKE_BLUE,
                    alignment=PP_ALIGN.CENTER)

    # Tagline
    if subtitle_text:
        add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                    subtitle_text, font_size=16, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

    # Stats line
    if tagline:
        add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                    tagline, font_size=14, color=MED_GRAY,
                    alignment=PP_ALIGN.CENTER)

    add_footer(slide, prs)
    return slide


def build_content_slide(prs, elements, is_demo_script=False):
    """Build a standard content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Blue accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 SLIDE_WIDTH, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SNOWFLAKE_BLUE
    bar.line.fill.background()

    y = Inches(0.4)
    left_margin = Inches(0.7)
    content_width = Inches(11.9)
    speaker_notes = []

    for etype, content in elements:
        if etype == "h1":
            add_textbox(slide, left_margin, y, content_width, Inches(0.6),
                        content, font_size=32, bold=True, color=DARK_BG)
            y += Inches(0.65)

        elif etype == "h2":
            add_textbox(slide, left_margin, y, content_width, Inches(0.5),
                        content, font_size=22, bold=False, color=SNOWFLAKE_BLUE)
            y += Inches(0.55)

        elif etype == "h3":
            add_textbox(slide, left_margin, y, content_width, Inches(0.4),
                        content, font_size=18, bold=True, color=DARK_BG)
            y += Inches(0.45)

        elif etype == "text":
            clean = content.replace("**", "")
            add_textbox(slide, left_margin, y, content_width, Inches(0.35),
                        clean, font_size=14, color=DARK_TEXT)
            y += Inches(0.38)

        elif etype == "bullets":
            remaining = SLIDE_HEIGHT - y - Inches(0.6)
            tf = add_rich_bullets(slide, left_margin, y, content_width,
                                  remaining, content, font_size=13)
            # Estimate height
            n_lines = len(content)
            est_height = max(Inches(0.3 * n_lines), Inches(0.6))
            y += est_height + Inches(0.15)

        elif etype == "table":
            rows = parse_md_table(content)
            if rows:
                remaining_h = SLIDE_HEIGHT - y - Inches(0.6)
                # Auto-size: fewer rows → bigger row height
                row_count = len(rows)
                font_size = 11 if row_count <= 10 else 9
                add_table(slide, left_margin, y, content_width, rows)
                y += Inches(0.33 * row_count) + Inches(0.15)

        elif etype == "code":
            code_h = add_code_block(slide, left_margin, y, content_width,
                                    content, font_size=9)
            y += code_h + Inches(0.15)

        elif etype == "quote":
            # Quotes go into speaker notes for demo scripts
            speaker_notes.append(content)
            # Also render as italic text on slide
            tf = add_textbox(slide, Inches(1.2), y, Inches(10.5), Inches(0.35),
                             f'"{content}"', font_size=13, color=SNOWFLAKE_BLUE)
            tf.paragraphs[0].font.italic = True
            y += Inches(0.4)

    # Add accumulated speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "\n".join(speaker_notes)

    add_footer(slide, prs)
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    md_text = SRC.read_text()

    # Split on --- separator
    raw_slides = re.split(r"\n---\n", md_text)

    # Filter out empty / comment-only blocks
    slide_blocks = []
    for block in raw_slides:
        cleaned = block.strip()
        # Skip empty or comment-only blocks
        if not cleaned or re.match(r"^<!--.*-->$", cleaned, re.DOTALL):
            continue
        # Skip blocks that are only HTML comments
        without_comments = re.sub(r"<!--.*?-->", "", cleaned, flags=re.DOTALL).strip()
        if not without_comments:
            continue
        slide_blocks.append(cleaned)

    print(f"Found {len(slide_blocks)} slide blocks")

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for idx, block in enumerate(slide_blocks):
        elements = split_slide_content(block)
        if not elements:
            continue

        is_title = idx == 0
        is_demo = any("Demo Script" in (c if isinstance(c, str) else "")
                       for _, c in elements if isinstance(c, str))

        if is_title:
            build_title_slide(prs, elements)
        else:
            build_content_slide(prs, elements, is_demo_script=is_demo)

        # Get slide title for logging
        title = ""
        for etype, content in elements:
            if etype == "h1":
                title = content
                break
        print(f"  Slide {idx + 1}: {title[:60]}")

    prs.save(str(DST))
    print(f"\nSaved: {DST}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
